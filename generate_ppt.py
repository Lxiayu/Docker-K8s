#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x00, 0x66, 0xFF)
CYAN = RGBColor(0x00, 0xFF, 0xFF)
DARK = RGBColor(0x0A, 0x0A, 0x0F)
DARK_CARD = RGBColor(0x12, 0x12, 0x1A)
DARK_SURFACE = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_100 = RGBColor(0xF0, 0xF0, 0xF5)
GRAY_300 = RGBColor(0xB0, 0xB0, 0xC0)
GRAY_500 = RGBColor(0x6B, 0x6B, 0x80)
GREEN = RGBColor(0x00, 0xFF, 0x88)
ORANGE = RGBColor(0xFF, 0x88, 0x00)
RED = RGBColor(0xFF, 0x33, 0x66)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)


def set_slide_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color if fill_color else DARK_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color if fill_color else DARK_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(tf, text, font_size=14, color=GRAY_300, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4), font_name='Microsoft YaHei'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_gradient_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total=12):
    add_text_box(slide, Inches(12.0), Inches(0.3), Inches(1.2), Inches(0.4),
                 f"{num:02d} / {total}", font_size=11, color=GRAY_500, alignment=PP_ALIGN.RIGHT,
                 font_name='Consolas')


def add_section_title(slide, title, subtitle=""):
    add_gradient_bar(slide, Inches(0.8), Inches(0.9), Inches(0.8), Pt(4))
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
                 title, font_size=36, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.75), Inches(10), Inches(0.5),
                     subtitle, font_size=16, color=GRAY_300)


def add_card(slide, left, top, width, height, title, desc, icon="", title_color=CYAN, border_color=None):
    card = add_shape(slide, left, top, width, height, fill_color=DARK_CARD, border_color=border_color)
    y_offset = top + Inches(0.2)
    if icon:
        add_text_box(slide, left + Inches(0.25), y_offset, Inches(0.5), Inches(0.5),
                     icon, font_size=24, color=title_color)
        y_offset += Inches(0.45)
    add_text_box(slide, left + Inches(0.25), y_offset, width - Inches(0.5), Inches(0.4),
                 title, font_size=16, color=title_color, bold=True)
    add_text_box(slide, left + Inches(0.25), y_offset + Inches(0.35), width - Inches(0.5), height - Inches(1.0),
                 desc, font_size=12, color=GRAY_300)
    return card


# ==================== SLIDE 1: COVER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_gradient_bar(slide, Inches(0), Inches(0), prs.slide_width, Pt(4))

add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.5),
             "CLOUD NATIVE  |  CI/CD  |  KUBERNETES", font_size=14, color=CYAN, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
             "云原生容器化\nCI/CD 自动化部署系统", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
             "代码即部署，部署即生效 —— 从代码提交到生产部署的全流程自动化", font_size=18, color=GRAY_300, alignment=PP_ALIGN.CENTER)

tags = ["Go + Gin", "React 18", "Docker + K8s", "Prometheus + Grafana", "Harbor + ELK"]
tag_x = Inches(2.5)
for tag in tags:
    shape = add_shape(slide, tag_x, Inches(4.8), Inches(1.5), Inches(0.4),
                      fill_color=RGBColor(0x0D, 0x1B, 0x3E), border_color=BLUE)
    shape.text_frame.paragraphs[0].text = tag
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = CYAN
    shape.text_frame.paragraphs[0].font.name = 'Consolas'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tag_x += Inches(1.7)

add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4),
             "答辩讲解", font_size=14, color=GRAY_500, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1)

# ==================== SLIDE 2: BACKGROUND ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "项目背景与痛点", "传统部署方式在微服务和云原生时代面临严峻挑战")
add_slide_number(slide, 2)

pain_points = [
    ("😰", "环境不一致", "开发、测试、生产环境配置差异\n导致部署失败率高，排查成本高", RED),
    ("🐢", "效率低下", "人工部署耗时长，发布频率受限\n业务迭代速度慢", ORANGE),
    ("💥", "故障恢复慢", "出现问题后手动回滚耗时长\n业务损失扩大，MTTR难以保障", PURPLE),
    ("🔒", "安全隐患", "依赖冲突、镜像漏洞、权限混乱\n缺乏统一安全管控", BLUE),
]

for i, (icon, title, desc, color) in enumerate(pain_points):
    col = i % 2
    row = i // 2
    left = Inches(0.8) + col * Inches(5.9)
    top = Inches(2.6) + row * Inches(1.85)
    card = add_shape(slide, left, top, Inches(5.6), Inches(1.65),
                     fill_color=DARK_CARD, border_color=RGBColor(0x25, 0x25, 0x35))
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.4),
                 icon, font_size=22, color=color)
    add_text_box(slide, left + Inches(0.7), top + Inches(0.15), Inches(4.5), Inches(0.35),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, left + Inches(0.7), top + Inches(0.55), Inches(4.5), Inches(0.9),
                 desc, font_size=12, color=GRAY_300)

stats_box = add_shape(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.8),
                      fill_color=RGBColor(0x0D, 0x1B, 0x3E), border_color=RGBColor(0x1A, 0x3A, 0x6E))
add_text_box(slide, Inches(1.2), Inches(6.5), Inches(3), Inches(0.5),
             "85% 企业因依赖升级导致中断", font_size=13, color=CYAN, bold=True)
add_text_box(slide, Inches(5.0), Inches(6.5), Inches(3), Inches(0.5),
             "70% 人工干预错误可被自动化消除", font_size=13, color=CYAN, bold=True)
add_text_box(slide, Inches(8.8), Inches(6.5), Inches(3), Inches(0.5),
             "80%+ 部署时间可被缩短", font_size=13, color=CYAN, bold=True)

# ==================== SLIDE 3: SYSTEM OVERVIEW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "系统概述", "一套完整的云原生容器化 CI/CD 自动化部署系统")
add_slide_number(slide, 3)

features = [
    ("🔄", "全流程自动化", "代码提交→镜像构建→测试验证→自动部署\n端到端零人工干预", BLUE),
    ("🌐", "多环境支持", "开发、测试、生产环境隔离\n与统一管理，配置即代码", CYAN),
    ("🚀", "高级部署策略", "滚动更新、灰度发布、蓝绿部署\n支持自动回滚", GREEN),
    ("📊", "实时监控告警", "Prometheus + Grafana\n全方位监控，多渠道告警", ORANGE),
    ("🛡️", "安全管理", "镜像漏洞扫描、RBAC权限控制\n全量审计日志", PURPLE),
    ("💻", "Web可视化", "现代化前端界面，操作简单直观\n支持明暗主题切换", RED),
]

for i, (icon, title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.8) + col * Inches(4.0)
    top = Inches(2.5) + row * Inches(2.3)
    add_card(slide, left, top, Inches(3.7), Inches(2.0), title, desc, icon, color)

# ==================== SLIDE 4: ARCHITECTURE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "系统架构设计", "三层分层架构：用户界面层 → 应用服务层 → 基础设施层")
add_slide_number(slide, 4)

layers = [
    ("🖥️  用户界面层  UI LAYER", [
        "React 18 + TypeScript", "Radix UI + Tailwind CSS",
        "Zustand 状态管理", "Vite 构建", "Recharts 图表"
    ], RGBColor(0x0D, 0x1B, 0x3E), BLUE),
    ("⚙️  应用服务层  SERVICE LAYER", [
        "Go 1.21 + Gin", "GORM + PostgreSQL",
        "JWT 认证", "Zap 日志", "Viper 配置", "Redis 缓存"
    ], RGBColor(0x16, 0x0D, 0x2E), PURPLE),
    ("🏗️  基础设施层  INFRA LAYER", [
        "Docker", "Kubernetes", "Harbor",
        "Prometheus + Grafana", "ELK Stack", "Vault"
    ], RGBColor(0x0D, 0x1E, 0x14), GREEN),
]

for i, (title, items, bg_color, title_color) in enumerate(layers):
    top = Inches(2.5) + i * Inches(1.65)
    layer = add_shape(slide, Inches(0.8), top, Inches(11.7), Inches(1.45),
                      fill_color=bg_color, border_color=title_color)
    add_text_box(slide, Inches(1.1), top + Inches(0.1), Inches(5), Inches(0.4),
                 title, font_size=15, color=title_color, bold=True)
    item_x = Inches(1.1)
    for item in items:
        tag_shape = add_shape(slide, item_x, top + Inches(0.6), Inches(1.7), Inches(0.35),
                              fill_color=RGBColor(0x08, 0x08, 0x10), border_color=RGBColor(0x30, 0x30, 0x45))
        tag_shape.text_frame.paragraphs[0].text = item
        tag_shape.text_frame.paragraphs[0].font.size = Pt(10)
        tag_shape.text_frame.paragraphs[0].font.color.rgb = GRAY_300
        tag_shape.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'
        tag_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        item_x += Inches(1.85)

    if i < 2:
        add_text_box(slide, Inches(6.3), top + Inches(1.35), Inches(0.8), Inches(0.3),
                     "▼", font_size=18, color=title_color, alignment=PP_ALIGN.CENTER)

# ==================== SLIDE 5: CORE FLOW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "核心业务流程", "从代码提交到生产部署的完整 CI/CD 数据流")
add_slide_number(slide, 5)

flow_steps_1 = [
    ("📝", "代码提交"),
    ("🔗", "Webhook触发"),
    ("🔨", "CI流水线"),
    ("🐳", "镜像构建"),
    ("🔍", "镜像扫描"),
]

flow_steps_2 = [
    ("🏛️", "推送仓库"),
    ("☸️", "K8s部署"),
    ("✅", "测试验证"),
    ("🚀", "灰度发布"),
    ("📊", "监控告警"),
]

for row, steps in enumerate([flow_steps_1, flow_steps_2]):
    y = Inches(2.5) + row * Inches(1.6)
    x = Inches(0.5)
    for j, (icon, label) in enumerate(steps):
        step = add_shape(slide, x, y, Inches(1.8), Inches(1.1),
                         fill_color=DARK_SURFACE, border_color=BLUE)
        tf = step.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = icon
        tf.paragraphs[0].font.size = Pt(22)
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_300
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        x += Inches(2.0)
        if j < len(steps) - 1:
            add_text_box(slide, x - Inches(0.35), y + Inches(0.3), Inches(0.4), Inches(0.4),
                         "→", font_size=20, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

metrics = [
    ("⚡ 触发延迟 ≤ 5秒", "代码提交后 Webhook 实时触发流水线"),
    ("🔨 构建时间 ≤ 10分钟", "多阶段构建 + 缓存优化加速"),
    ("🚀 部署时间 ≤ 2分钟", "Kubernetes 滚动更新零停机"),
]

for i, (title, desc) in enumerate(metrics):
    left = Inches(0.8) + i * Inches(4.0)
    box = add_shape(slide, left, Inches(5.8), Inches(3.7), Inches(1.2),
                    fill_color=RGBColor(0x0D, 0x1B, 0x3E), border_color=RGBColor(0x1A, 0x3A, 0x6E))
    add_text_box(slide, left + Inches(0.2), Inches(5.9), Inches(3.3), Inches(0.4),
                 title, font_size=14, color=CYAN, bold=True)
    add_text_box(slide, left + Inches(0.2), Inches(6.3), Inches(3.3), Inches(0.4),
                 desc, font_size=11, color=GRAY_300)

# ==================== SLIDE 6: BACKEND ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "后端服务架构", "基于 Go + Gin 构建的高性能微服务后端，采用经典分层架构")
add_slide_number(slide, 6)

code_box = add_shape(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.8),
                     fill_color=RGBColor(0x08, 0x08, 0x10), border_color=RGBColor(0x25, 0x25, 0x35))
code_text = """// 请求处理链路
Router → Middleware → Handler → Service → Model

// 路由层 - API 定义与分组
api := engine.Group("/api/v1")
  ├── auth      → 登录 / 注册
  ├── users     → 用户管理 (JWT Auth)
  ├── pipelines → 流水线管理
  ├── deployments → 部署管理
  ├── images    → 镜像管理
  ├── monitoring → 监控告警
  ├── settings  → 系统配置
  ├── dashboard → 数据看板
  └── audit     → 审计日志"""
add_text_box(slide, Inches(1.0), Inches(2.6), Inches(5.1), Inches(3.6),
             code_text, font_size=11, color=GRAY_300, font_name='Consolas')

services = [
    ("AuthService", "JWT认证 · bcrypt加密", BLUE),
    ("PipelineService", "流水线编排 · 执行引擎", CYAN),
    ("K8sService", "集群管理 · 自动回滚", GREEN),
    ("BuildService", "多阶段构建 · 缓存优化", ORANGE),
    ("HarborService", "镜像仓库 · 漏洞扫描", PURPLE),
    ("RBACService", "角色权限 · 细粒度控制", RED),
    ("GitService", "Webhook · 仓库集成", BLUE),
    ("MonitoringService", "指标采集 · 告警规则", CYAN),
]

for i, (name, desc, color) in enumerate(services):
    col = i % 2
    row = i // 2
    left = Inches(6.8) + col * Inches(3.1)
    top = Inches(2.5) + row * Inches(1.0)
    card = add_shape(slide, left, top, Inches(2.9), Inches(0.85),
                     fill_color=DARK_CARD, border_color=RGBColor(0x25, 0x25, 0x35))
    add_text_box(slide, left + Inches(0.15), top + Inches(0.08), Inches(2.6), Inches(0.35),
                 name, font_size=13, color=color, bold=True, font_name='Consolas')
    add_text_box(slide, left + Inches(0.15), top + Inches(0.42), Inches(2.6), Inches(0.3),
                 desc, font_size=10, color=GRAY_300)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(2), Inches(0.3),
             "🛡️ 中间件链", font_size=14, color=CYAN, bold=True)

middleware_tags = ["JWT Auth", "CORS", "Logger", "Recovery", "RateLimit", "Performance"]
tag_x = Inches(2.8)
for tag in middleware_tags:
    tag_shape = add_shape(slide, tag_x, Inches(6.5), Inches(1.3), Inches(0.35),
                          fill_color=RGBColor(0x0D, 0x1B, 0x3E), border_color=BLUE)
    tag_shape.text_frame.paragraphs[0].text = tag
    tag_shape.text_frame.paragraphs[0].font.size = Pt(10)
    tag_shape.text_frame.paragraphs[0].font.color.rgb = CYAN
    tag_shape.text_frame.paragraphs[0].font.name = 'Consolas'
    tag_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tag_x += Inches(1.45)

# ==================== SLIDE 7: FRONTEND ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "前端控制台", "React 18 + TypeScript 构建的现代化 Web 控制台")
add_slide_number(slide, 7)

mock_browser = add_shape(slide, Inches(0.8), Inches(2.5), Inches(6.0), Inches(4.5),
                         fill_color=DARK_SURFACE, border_color=RGBColor(0x25, 0x25, 0x35))
toolbar = add_rect(slide, Inches(0.8), Inches(2.5), Inches(6.0), Inches(0.45),
                   fill_color=RGBColor(0x08, 0x08, 0x10))
for j, c in enumerate([RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFF, 0xBD, 0x2E), RGBColor(0x28, 0xC8, 0x40)]):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0) + j * Inches(0.25), Inches(2.6), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = c
    dot.line.fill.background()

add_text_box(slide, Inches(1.9), Inches(2.55), Inches(3), Inches(0.3),
             "localhost:3000/dashboard", font_size=10, color=GRAY_500, font_name='Consolas')

stat_cards = [
    ("128", "流水线总数", BLUE),
    ("96.5%", "构建成功率", GREEN),
    ("45", "活跃部署", PURPLE),
]
for i, (num, label, color) in enumerate(stat_cards):
    left = Inches(1.0) + i * Inches(1.8)
    stat_box = add_shape(slide, left, Inches(3.2), Inches(1.6), Inches(0.9),
                         fill_color=RGBColor(0x08, 0x08, 0x10), border_color=RGBColor(0x20, 0x20, 0x30))
    add_text_box(slide, left + Inches(0.1), Inches(3.25), Inches(1.4), Inches(0.45),
                 num, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(3.7), Inches(1.4), Inches(0.3),
                 label, font_size=10, color=GRAY_300, alignment=PP_ALIGN.CENTER)

builds_box = add_shape(slide, Inches(1.0), Inches(4.3), Inches(5.6), Inches(2.4),
                       fill_color=RGBColor(0x08, 0x08, 0x10), border_color=RGBColor(0x20, 0x20, 0x30))
add_text_box(slide, Inches(1.2), Inches(4.4), Inches(3), Inches(0.3),
             "最近构建", font_size=11, color=GRAY_500)
build_items = [
    ("●", "frontend-app #256", "2m ago", GREEN),
    ("●", "backend-api #189", "5m ago", GREEN),
    ("●", "auth-service #92", "12m ago", RED),
]
for i, (dot, name, time, color) in enumerate(build_items):
    y = Inches(4.75) + i * Inches(0.45)
    add_text_box(slide, Inches(1.3), y, Inches(0.2), Inches(0.3), dot, font_size=10, color=color)
    add_text_box(slide, Inches(1.6), y, Inches(2.5), Inches(0.3), name, font_size=12, color=WHITE)
    add_text_box(slide, Inches(4.5), y, Inches(1.5), Inches(0.3), time, font_size=10, color=GRAY_500, alignment=PP_ALIGN.RIGHT)

pages = [
    ("📊", "Dashboard 仪表盘", "系统概览、关键指标、实时状态"),
    ("🔄", "Pipelines 流水线", "创建、编辑、触发、查看构建日志"),
    ("🚀", "Deployments 部署管理", "多策略部署、回滚、Pod监控"),
    ("🐳", "Images 镜像管理", "镜像列表、构建、漏洞扫描"),
    ("📈", "Monitoring 监控告警", "实时指标、告警规则、通知配置"),
    ("👥", "Users 用户管理", "RBAC角色、权限分配、审计日志"),
]

for i, (icon, title, desc) in enumerate(pages):
    top = Inches(2.5) + i * Inches(0.8)
    card = add_shape(slide, Inches(7.2), top, Inches(5.3), Inches(0.7),
                     fill_color=DARK_CARD, border_color=RGBColor(0x25, 0x25, 0x35))
    add_text_box(slide, Inches(7.4), top + Inches(0.08), Inches(0.4), Inches(0.3),
                 icon, font_size=16, color=CYAN)
    add_text_box(slide, Inches(7.9), top + Inches(0.05), Inches(2.2), Inches(0.3),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.9), top + Inches(0.35), Inches(4.2), Inches(0.3),
                 desc, font_size=10, color=GRAY_300)

# ==================== SLIDE 8: DATABASE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "数据库设计", "PostgreSQL 数据库，GORM 自动迁移，完善的索引与约束设计")
add_slide_number(slide, 8)

table_data = [
    ("users", "用户表", "username, email, role, status", BLUE),
    ("roles", "角色表", "name, description", BLUE),
    ("permissions", "权限表", "resource, action", BLUE),
    ("git_repositories", "代码仓库", "url, type, branch, webhook", CYAN),
    ("pipelines", "流水线", "repo_id, config, status", CYAN),
    ("builds", "构建记录", "pipeline_id, commit_hash, duration", CYAN),
    ("deployments", "部署记录", "namespace, image, strategy, env", GREEN),
    ("audit_logs", "审计日志", "user_id, action, resource, ip", ORANGE),
]

table_top = Inches(2.5)
header = add_rect(slide, Inches(0.8), table_top, Inches(6.0), Inches(0.45),
                  fill_color=RGBColor(0x0D, 0x1B, 0x3E))
add_text_box(slide, Inches(1.0), table_top + Inches(0.05), Inches(1.5), Inches(0.35),
             "数据表", font_size=12, color=CYAN, bold=True)
add_text_box(slide, Inches(2.8), table_top + Inches(0.05), Inches(1.2), Inches(0.35),
             "说明", font_size=12, color=CYAN, bold=True)
add_text_box(slide, Inches(4.2), table_top + Inches(0.05), Inches(2.5), Inches(0.35),
             "关键字段", font_size=12, color=CYAN, bold=True)

for i, (table, desc, fields, color) in enumerate(table_data):
    y = table_top + Inches(0.45) + i * Inches(0.42)
    row_bg = DARK_CARD if i % 2 == 0 else RGBColor(0x0E, 0x0E, 0x16)
    add_rect(slide, Inches(0.8), y, Inches(6.0), Inches(0.42), fill_color=row_bg)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(1.5), Inches(0.3),
                 table, font_size=11, color=color, bold=True, font_name='Consolas')
    add_text_box(slide, Inches(2.8), y + Inches(0.05), Inches(1.2), Inches(0.3),
                 desc, font_size=11, color=GRAY_300)
    add_text_box(slide, Inches(4.2), y + Inches(0.05), Inches(2.5), Inches(0.3),
                 fields, font_size=10, color=GRAY_500, font_name='Consolas')

rbac_box = add_shape(slide, Inches(7.2), Inches(2.5), Inches(5.3), Inches(2.8),
                     fill_color=RGBColor(0x08, 0x08, 0x10), border_color=RGBColor(0x25, 0x25, 0x35))
rbac_text = """// RBAC 权限模型
Role ←(many2many)→ Permission

// 四级角色体系
┌──────────┬───────────────────────────┐
│ admin    │ 全部权限                   │
│ devops   │ 流水线+部署+仓库+监控       │
│ developer│ 只读权限                   │
│ viewer   │ 只读权限                   │
└──────────┴───────────────────────────┘

// 权限粒度: resource + action
pipelines:create  pipelines:trigger
deployments:execute  monitoring:update"""
add_text_box(slide, Inches(7.4), Inches(2.6), Inches(4.9), Inches(2.6),
             rbac_text, font_size=10, color=GRAY_300, font_name='Consolas')

highlights = [
    "✓ 软删除 (deleted_at) 保护数据",
    "✓ 自动更新时间触发器",
    "✓ 枚举类型约束 (build_status, deployment_strategy)",
    "✓ 审计日志自动清理函数 (90天)",
]
for i, h in enumerate(highlights):
    add_text_box(slide, Inches(7.4), Inches(5.5) + i * Inches(0.35), Inches(5), Inches(0.3),
                 h, font_size=12, color=GREEN)

# ==================== SLIDE 9: DEPLOYMENT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "部署策略与 Kubernetes 集成", "三种高级部署策略，保障零停机发布与快速故障恢复")
add_slide_number(slide, 9)

strategies = [
    ("🔄", "滚动更新", "Rolling Update", "逐步替换旧版本 Pod\n配置 maxSurge / maxUnavailable\n就绪探针保障可用性\n可用性 ≥ 99.99%", BLUE),
    ("🐦", "灰度发布", "Canary", "金丝雀发布策略\n流量权重 1% → 100%\n基于指标自动判断\n用户无感知切换", PURPLE),
    ("🔵🟢", "蓝绿部署", "Blue-Green", "两套完整环境切换\n零风险即时回滚\n切换流量即可完成\n回滚时间 ≤ 2分钟", CYAN),
]

for i, (icon, name, en_name, desc, color) in enumerate(strategies):
    left = Inches(0.8) + i * Inches(4.1)
    card = add_shape(slide, left, Inches(2.5), Inches(3.8), Inches(2.5),
                     fill_color=DARK_CARD, border_color=color)
    add_rect(slide, left, Inches(2.5), Inches(3.8), Pt(4), fill_color=color)
    add_text_box(slide, left + Inches(0.2), Inches(2.7), Inches(3.4), Inches(0.5),
                 icon + "  " + name, font_size=20, color=color, bold=True)
    add_text_box(slide, left + Inches(0.2), Inches(3.2), Inches(3.4), Inches(0.3),
                 en_name, font_size=11, color=GRAY_500, font_name='Consolas')
    add_text_box(slide, left + Inches(0.2), Inches(3.6), Inches(3.4), Inches(1.2),
                 desc, font_size=12, color=GRAY_300)

k8s_box = add_shape(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(1.8),
                    fill_color=RGBColor(0x08, 0x08, 0x10), border_color=RGBColor(0x25, 0x25, 0x35))
k8s_text = """// Kubernetes 部署架构
├── Deployment   → 应用部署
├── Service      → 服务暴露
├── Ingress      → 流量路由
├── HPA          → 弹性伸缩
├── PDB          → Pod中断预算
├── ConfigMap    → 配置管理
├── Secret       → 密钥管理
├── NetworkPolicy→ 网络隔离
└── PVC          → 持久化存储"""
add_text_box(slide, Inches(1.0), Inches(5.35), Inches(5.1), Inches(1.7),
             k8s_text, font_size=10, color=GRAY_300, font_name='Consolas')

ha_items = [
    ("多副本部署", "前端/后端多副本 + PDB 保障最小可用实例"),
    ("自动回滚机制", "监控指标异常 → 自动触发回滚 → 切换稳定版本"),
    ("网络策略隔离", "默认拒绝入站 + 数据库/Redis 专用网络策略"),
]
for i, (title, desc) in enumerate(ha_items):
    top = Inches(5.3) + i * Inches(0.6)
    box = add_shape(slide, Inches(6.8), top, Inches(5.7), Inches(0.55),
                    fill_color=RGBColor(0x0D, 0x1B, 0x3E), border_color=RGBColor(0x1A, 0x3A, 0x6E))
    add_text_box(slide, Inches(7.0), top + Inches(0.03), Inches(2), Inches(0.25),
                 title, font_size=12, color=CYAN, bold=True)
    add_text_box(slide, Inches(7.0), top + Inches(0.28), Inches(5.2), Inches(0.25),
                 desc, font_size=10, color=GRAY_300)

# ==================== SLIDE 10: MONITORING ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "监控告警体系", "Prometheus + Grafana + AlertManager + ELK 全方位可观测性")
add_slide_number(slide, 10)

monitor_flow = [
    ("📡", "Prometheus\n采集"),
    ("💾", "时序数据\n存储"),
    ("📊", "Grafana\n可视化"),
    ("🔔", "AlertManager\n告警"),
]

for i, (icon, label) in enumerate(monitor_flow):
    left = Inches(0.8) + i * Inches(2.8)
    step = add_shape(slide, left, Inches(2.5), Inches(2.2), Inches(1.0),
                     fill_color=DARK_SURFACE, border_color=BLUE)
    tf = step.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].text = icon
    tf.paragraphs[0].font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_300
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    if i < 3:
        add_text_box(slide, left + Inches(2.25), Inches(2.75), Inches(0.5), Inches(0.4),
                     "→", font_size=18, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

alert_header = add_rect(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.4),
                        fill_color=RGBColor(0x0D, 0x1B, 0x3E))
add_text_box(slide, Inches(1.0), Inches(3.83), Inches(1.2), Inches(0.3),
             "类别", font_size=11, color=CYAN, bold=True)
add_text_box(slide, Inches(2.3), Inches(3.83), Inches(2.5), Inches(0.3),
             "规则", font_size=11, color=CYAN, bold=True)
add_text_box(slide, Inches(5.0), Inches(3.83), Inches(1), Inches(0.3),
             "级别", font_size=11, color=CYAN, bold=True)

alerts = [
    ("节点", "CPU > 80% / 内存 > 85%", "P0", RED),
    ("Pod", "CrashLoopBackOff / Pending", "P0", RED),
    ("服务", "错误率 > 5% / 响应 > 2s", "P1", ORANGE),
    ("资源", "磁盘 > 90% / PVC 即将满", "P1", ORANGE),
    ("业务", "部署失败 / 构建超时", "P2", PURPLE),
]

for i, (cat, rule, level, color) in enumerate(alerts):
    y = Inches(4.2) + i * Inches(0.38)
    row_bg = DARK_CARD if i % 2 == 0 else RGBColor(0x0E, 0x0E, 0x16)
    add_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.38), fill_color=row_bg)
    add_text_box(slide, Inches(1.0), y + Inches(0.03), Inches(1.2), Inches(0.3),
                 cat, font_size=11, color=GRAY_300)
    add_text_box(slide, Inches(2.3), y + Inches(0.03), Inches(2.5), Inches(0.3),
                 rule, font_size=10, color=GRAY_300)
    level_shape = add_shape(slide, Inches(5.0), y + Inches(0.04), Inches(0.5), Inches(0.28),
                            fill_color=RGBColor(0x08, 0x08, 0x10), border_color=color)
    level_shape.text_frame.paragraphs[0].text = level
    level_shape.text_frame.paragraphs[0].font.size = Pt(9)
    level_shape.text_frame.paragraphs[0].font.color.rgb = color
    level_shape.text_frame.paragraphs[0].font.bold = True
    level_shape.text_frame.paragraphs[0].font.name = 'Consolas'
    level_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

notify_channels = [
    ("📧", "邮件通知", "SMTP 邮件告警，支持自定义模板"),
    ("💬", "钉钉机器人", "Webhook 推送，实时告警到群组"),
    ("🏢", "企业微信", "企业微信应用消息推送"),
]
for i, (icon, title, desc) in enumerate(notify_channels):
    top = Inches(3.8) + i * Inches(0.85)
    card = add_shape(slide, Inches(6.8), top, Inches(5.7), Inches(0.75),
                     fill_color=DARK_CARD, border_color=RGBColor(0x25, 0x25, 0x35))
    add_text_box(slide, Inches(7.0), top + Inches(0.08), Inches(0.4), Inches(0.3),
                 icon, font_size=18, color=CYAN)
    add_text_box(slide, Inches(7.5), top + Inches(0.05), Inches(2), Inches(0.3),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.5), top + Inches(0.38), Inches(4.5), Inches(0.3),
                 desc, font_size=10, color=GRAY_300)

elk_box = add_shape(slide, Inches(6.8), Inches(6.4), Inches(5.7), Inches(0.8),
                    fill_color=RGBColor(0x08, 0x08, 0x10), border_color=RGBColor(0x25, 0x25, 0x35))
add_text_box(slide, Inches(7.0), Inches(6.45), Inches(5.3), Inches(0.3),
             "📋 ELK 日志: Application → Fluent Bit → Elasticsearch → Kibana",
             font_size=11, color=GRAY_300, font_name='Consolas')
add_text_box(slide, Inches(7.0), Inches(6.8), Inches(5.3), Inches(0.3),
             "结构化JSON · ILM生命周期 · 日志90天 · 审计365天", font_size=10, color=GRAY_500)

# ==================== SLIDE 11: SECURITY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "安全设计", "多层次安全防护，从传输加密到审计追踪的全方位安全体系")
add_slide_number(slide, 11)

sec_cards = [
    ("🔑", "身份认证", "JWT Token\nbcrypt 密码加密\n登录限流保护", BLUE),
    ("🛡️", "权限控制", "RBAC 角色模型\n细粒度权限\n最小权限原则", PURPLE),
    ("🔐", "传输加密", "TLS 1.3\nHTTPS 全链路\nAES-256 存储", GREEN),
    ("🔍", "镜像安全", "Trivy 漏洞扫描\n高危阻断策略\nHarbor 访问控制", ORANGE),
]

for i, (icon, title, desc, color) in enumerate(sec_cards):
    left = Inches(0.8) + i * Inches(3.1)
    add_card(slide, left, Inches(2.5), Inches(2.8), Inches(1.8), title, desc, icon, color)

net_box = add_shape(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.5),
                    fill_color=RGBColor(0x08, 0x08, 0x10), border_color=RGBColor(0x25, 0x25, 0x35))
net_text = """// Kubernetes NetworkPolicy

// 默认拒绝所有入站
default-deny → 拒绝所有入站流量

// 数据库专用策略
allow-database → 仅允许后端服务访问

// Redis 专用策略
allow-redis → 仅允许后端服务访问

// Ingress 控制
Nginx Ingress → 限制外部访问入口"""
add_text_box(slide, Inches(1.0), Inches(4.7), Inches(5.1), Inches(2.3),
             net_text, font_size=10, color=GRAY_300, font_name='Consolas')

audit_box = add_shape(slide, Inches(6.8), Inches(4.6), Inches(5.7), Inches(2.5),
                      fill_color=RGBColor(0x08, 0x08, 0x10), border_color=RGBColor(0x25, 0x25, 0x35))
audit_text = """// 审计日志记录内容
{
  "user_id":   uint,
  "username":  string,
  "action":    string,    // 操作类型
  "resource":  string,    // 操作资源
  "method":    string,    // HTTP方法
  "path":      string,    // 请求路径
  "ip":        string,    // 来源IP
  "status":    int        // 响应状态
}

// 自动清理: 保留90天
cleanup_old_audit_logs(90)"""
add_text_box(slide, Inches(7.0), Inches(4.7), Inches(5.3), Inches(2.3),
             audit_text, font_size=10, color=GRAY_300, font_name='Consolas')

# ==================== SLIDE 12: SUMMARY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 12)

add_gradient_bar(slide, Inches(5.5), Inches(0.9), Inches(2.3), Pt(4))
add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.8),
             "项目总结与展望", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

stat_items = [
    ("8+", "核心服务模块", BLUE),
    ("50+", "API 接口", PURPLE),
    ("3", "部署策略", GREEN),
]
for i, (num, label, color) in enumerate(stat_items):
    left = Inches(1.5) + i * Inches(3.5)
    card = add_shape(slide, left, Inches(2.0), Inches(3.0), Inches(1.2),
                     fill_color=DARK_CARD, border_color=color)
    add_rect(slide, left, Inches(2.0), Inches(3.0), Pt(3), fill_color=color)
    add_text_box(slide, left + Inches(0.2), Inches(2.15), Inches(2.6), Inches(0.6),
                 num, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(2.75), Inches(2.6), Inches(0.3),
                 label, font_size=13, color=GRAY_300, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(5), Inches(0.4),
             "🎯 项目亮点", font_size=16, color=CYAN, bold=True)

highlights = [
    "完整的 CI/CD 全流程自动化闭环",
    "云原生架构，K8s 弹性伸缩",
    "多策略部署 + 自动回滚保障",
    "全方位监控告警与日志体系",
    "RBAC + 审计日志安全体系",
    "现代化 React 前端控制台",
]
for i, h in enumerate(highlights):
    col = i % 2
    row = i // 2
    left = Inches(1.5) + col * Inches(5.2)
    top = Inches(4.0) + row * Inches(0.38)
    add_text_box(slide, left, top, Inches(5), Inches(0.3),
                 "✦  " + h, font_size=13, color=GRAY_300)

add_text_box(slide, Inches(1.5), Inches(5.3), Inches(5), Inches(0.4),
             "🚀 未来展望", font_size=16, color=CYAN, bold=True)

timeline = [
    ("阶段一 (已完成)", "基础设施搭建 · 后端核心服务 · 前端控制台 · 认证权限体系"),
    ("阶段二 (进行中)", "灰度发布完善 · 镜像漏洞扫描集成 · 多集群管理 · 插件系统"),
    ("阶段三 (规划中)", "多云/混合云部署 · GitOps (ArgoCD) · AI 智能运维 · 服务网格 (Istio)"),
]

for i, (phase, desc) in enumerate(timeline):
    top = Inches(5.8) + i * Inches(0.5)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.7), top + Inches(0.08), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BLUE
    dot.line.fill.background()
    if i < 2:
        line = add_rect(slide, Inches(1.75), top + Inches(0.23), Pt(2), Inches(0.35), fill_color=BLUE)
    add_text_box(slide, Inches(2.0), top, Inches(2.5), Inches(0.3),
                 phase, font_size=12, color=CYAN, bold=True)
    add_text_box(slide, Inches(4.5), top, Inches(7), Inches(0.3),
                 desc, font_size=11, color=GRAY_300)

add_text_box(slide, Inches(1.5), Inches(7.0), Inches(10), Inches(0.4),
             "感谢聆听  |  云原生容器化 CI/CD 自动化部署系统  |  代码即部署，部署即生效",
             font_size=14, color=GRAY_500, alignment=PP_ALIGN.CENTER)

# ==================== SAVE ====================
output_path = '/workspace/云原生CICD自动化部署系统_答辩讲解.pptx'
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
